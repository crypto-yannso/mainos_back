#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Module d'export pour la conversion des rapports Mainos dans différents formats.
Supporte l'export vers Markdown, PDF, PowerPoint, Word et HTML.
Fournit également des fonctionnalités pour l'intégration avec des services externes.
"""

import os
import re
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Union

from src.open_deep_research.mainos_types import FormatOutput, PlateformeIntegration

# Importations conditionnelles pour éviter des dépendances obligatoires
try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

try:
    from weasyprint import HTML
    WEASYPRINT_AVAILABLE = True
except ImportError:
    WEASYPRINT_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def exporter_rapport(
    contenu: str,
    format_sortie: FormatOutput,
    chemin_sortie: Optional[str] = None,
    metadata: Optional[Dict] = None,
) -> str:
    """
    Exporte un rapport dans le format spécifié.
    
    Args:
        contenu: Le contenu du rapport au format Markdown
        format_sortie: Le format de sortie souhaité
        chemin_sortie: Le chemin du fichier de sortie (optionnel)
        metadata: Métadonnées supplémentaires pour le document (optionnel)
        
    Returns:
        Le chemin du fichier exporté
    """
    if chemin_sortie is None:
        # Créer un nom de fichier temporaire si aucun n'est fourni
        extension = format_sortie.value.lower()
        fd, chemin_sortie = tempfile.mkstemp(suffix=f".{extension}")
        os.close(fd)
    
    if format_sortie == FormatOutput.MARKDOWN:
        # Simplement écrire le contenu Markdown
        with open(chemin_sortie, "w", encoding="utf-8") as f:
            f.write(contenu)
    
    elif format_sortie == FormatOutput.PDF:
        convertir_markdown_pdf(contenu, chemin_sortie, metadata)
    
    elif format_sortie == FormatOutput.HTML:
        convertir_markdown_html(contenu, chemin_sortie, metadata)
    
    elif format_sortie == FormatOutput.DOCX:
        convertir_markdown_docx(contenu, chemin_sortie, metadata)
    
    elif format_sortie == FormatOutput.PPTX:
        convertir_markdown_pptx(contenu, chemin_sortie, metadata)
    
    else:
        raise ValueError(f"Format de sortie non pris en charge : {format_sortie}")
    
    return chemin_sortie


def convertir_markdown_html(
    markdown_text: str, 
    output_path: str, 
    metadata: Optional[Dict] = None
) -> str:
    """
    Convertit du texte Markdown en HTML et l'enregistre dans un fichier.
    
    Args:
        markdown_text: Le texte Markdown à convertir
        output_path: Le chemin du fichier HTML de sortie
        metadata: Métadonnées pour le document HTML (titre, auteur, etc.)
        
    Returns:
        Le chemin du fichier HTML généré
    """
    if not MARKDOWN_AVAILABLE:
        raise ImportError("Le module 'markdown' est requis pour la conversion HTML.")
    
    # Convertir le Markdown en HTML
    html_content = markdown.markdown(
        markdown_text,
        extensions=['extra', 'smarty', 'tables', 'toc']
    )
    
    # Créer un document HTML complet avec les métadonnées
    title = metadata.get("title", "Rapport Mainos") if metadata else "Rapport Mainos"
    full_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <style>
        body {{
            font-family: Arial, sans-serif;
            line-height: 1.6;
            max-width: 900px;
            margin: 0 auto;
            padding: 1em;
        }}
        h1, h2, h3, h4, h5, h6 {{
            color: #333;
            margin-top: 1.5em;
        }}
        table {{
            border-collapse: collapse;
            width: 100%;
        }}
        table, th, td {{
            border: 1px solid #ddd;
        }}
        th, td {{
            padding: 10px;
            text-align: left;
        }}
        img {{
            max-width: 100%;
        }}
        blockquote {{
            border-left: 5px solid #eee;
            padding-left: 1em;
            color: #555;
        }}
        pre {{
            background: #f5f5f5;
            padding: 1em;
            overflow-x: auto;
            border-radius: 4px;
        }}
        code {{
            background: #f5f5f5;
            padding: 2px 4px;
            border-radius: 4px;
        }}
    </style>
</head>
<body>
    {html_content}
</body>
</html>
"""
    
    # Écrire le HTML dans le fichier de sortie
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(full_html)
    
    return output_path


def convertir_markdown_pdf(
    markdown_text: str, 
    output_path: str, 
    metadata: Optional[Dict] = None
) -> str:
    """
    Convertit du texte Markdown en PDF et l'enregistre dans un fichier.
    
    Args:
        markdown_text: Le texte Markdown à convertir
        output_path: Le chemin du fichier PDF de sortie
        metadata: Métadonnées pour le document PDF
        
    Returns:
        Le chemin du fichier PDF généré
    """
    if not WEASYPRINT_AVAILABLE:
        raise ImportError("Le module 'weasyprint' est requis pour la conversion PDF.")
    
    # Convertir d'abord en HTML en utilisant la fonction existante
    with tempfile.NamedTemporaryFile(suffix=".html", delete=False) as temp_file:
        html_path = temp_file.name
    
    try:
        convertir_markdown_html(markdown_text, html_path, metadata)
        
        # Utiliser WeasyPrint pour convertir HTML en PDF
        html = HTML(filename=html_path)
        html.write_pdf(output_path)
    finally:
        # Nettoyer le fichier temporaire
        if os.path.exists(html_path):
            os.unlink(html_path)
    
    return output_path


def convertir_markdown_docx(
    markdown_text: str, 
    output_path: str, 
    metadata: Optional[Dict] = None
) -> str:
    """
    Convertit du texte Markdown en document Word et l'enregistre.
    
    Args:
        markdown_text: Le texte Markdown à convertir
        output_path: Le chemin du fichier DOCX de sortie
        metadata: Métadonnées pour le document
        
    Returns:
        Le chemin du fichier DOCX généré
    """
    if not DOCX_AVAILABLE:
        raise ImportError("Le module 'python-docx' est requis pour la conversion DOCX.")
    
    # Créer un nouveau document Word
    doc = Document()
    
    # Ajouter les métadonnées
    if metadata:
        doc.core_properties.title = metadata.get("title", "")
        doc.core_properties.author = metadata.get("author", "Mainos")
        doc.core_properties.category = metadata.get("category", "")
    
    # Parser le Markdown et le convertir en éléments Word
    # Nous utilisons une approche simplifiée ici
    
    # Diviser le texte en lignes
    lines = markdown_text.split('\n')
    
    # Variables pour suivre l'état du parsing
    in_list = False
    in_code_block = False
    current_list_items = []
    
    # Regex pour les titres, les listes, etc.
    heading_pattern = re.compile(r'^(#{1,6})\s+(.+)$')
    list_item_pattern = re.compile(r'^(\s*[-*+]|\s*\d+\.)\s+(.+)$')
    code_block_pattern = re.compile(r'^```\w*$')
    
    for line in lines:
        # Vérifier si nous sommes dans un bloc de code
        if code_block_pattern.match(line):
            in_code_block = not in_code_block
            if in_code_block:
                p = doc.add_paragraph()
                p.style = 'Code'
            continue
        
        # Si nous sommes dans un bloc de code, ajouter la ligne telle quelle
        if in_code_block:
            p = doc.add_paragraph(line)
            p.style = 'Code'
            continue
        
        # Traiter les titres
        heading_match = heading_pattern.match(line)
        if heading_match:
            level = len(heading_match.group(1))
            text = heading_match.group(2)
            
            # Ajouter les éléments de liste en attente avant le titre
            if in_list and current_list_items:
                for item in current_list_items:
                    doc.add_paragraph(item, style='ListBullet')
                current_list_items = []
                in_list = False
            
            # Ajouter le titre
            doc.add_heading(text, level=level)
            continue
        
        # Traiter les éléments de liste
        list_match = list_item_pattern.match(line)
        if list_match:
            item_text = list_match.group(2)
            current_list_items.append(item_text)
            in_list = True
            continue
        
        # Si nous étions dans une liste et maintenant une ligne non-liste
        if in_list and line.strip() == "":
            for item in current_list_items:
                doc.add_paragraph(item, style='ListBullet')
            current_list_items = []
            in_list = False
            doc.add_paragraph()  # Ajouter un paragraphe vide
            continue
        
        # Traiter les paragraphes normaux
        if line.strip():
            # Traiter les lignes non vides comme des paragraphes
            doc.add_paragraph(line)
        else:
            # Lignes vides -> saut de paragraphe
            if not in_list:
                doc.add_paragraph()
    
    # S'assurer que tous les éléments de liste sont ajoutés
    if in_list and current_list_items:
        for item in current_list_items:
            doc.add_paragraph(item, style='ListBullet')
    
    # Enregistrer le document Word
    doc.save(output_path)
    return output_path


def convertir_markdown_pptx(
    markdown_text: str, 
    output_path: str, 
    metadata: Optional[Dict] = None
) -> str:
    """
    Convertit du texte Markdown en présentation PowerPoint et l'enregistre.
    
    Args:
        markdown_text: Le texte Markdown à convertir
        output_path: Le chemin du fichier PPTX de sortie
        metadata: Métadonnées pour la présentation
        
    Returns:
        Le chemin du fichier PPTX généré
    """
    if not PPTX_AVAILABLE:
        raise ImportError("Le module 'python-pptx' est requis pour la conversion PPTX.")
    
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Obtenir le titre de la présentation depuis les métadonnées
    title = metadata.get("title", "Rapport Mainos") if metadata else "Rapport Mainos"
    
    # Ajouter une diapositive de titre
    slide_layout = prs.slide_layouts[0]  # Layout de titre
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    author = metadata.get("author", "Généré par Mainos") if metadata else "Généré par Mainos"
    subtitle_shape.text = author
    
    # Parser le Markdown pour créer des diapositives
    # Diviser par les titres de niveau 1 ou 2
    sections = re.split(r'(?m)^#{1,2}\s+', markdown_text)
    titles = re.findall(r'(?m)^#{1,2}\s+([^\n]+)', markdown_text)
    
    # La première section pourrait être une introduction avant le premier titre
    if sections[0].strip() and not sections[0].startswith('#'):
        add_content_slide(prs, "Introduction", sections[0])
    
    # Traiter chaque section avec son titre
    for i, (section_content, section_title) in enumerate(zip(sections[1:], titles)):
        add_content_slide(prs, section_title, section_content)
    
    # Enregistrer la présentation
    prs.save(output_path)
    return output_path


def add_content_slide(presentation, title, content):
    """
    Ajoute une diapositive avec titre et contenu à la présentation.
    
    Args:
        presentation: L'objet Presentation
        title: Le titre de la diapositive
        content: Le contenu à ajouter
    """
    # Utiliser un layout avec titre et contenu
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    
    # Définir le titre
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Ajouter le contenu
    content_shape = slide.placeholders[1]
    
    # Limiter le contenu pour qu'il tienne sur une diapositive
    # Prendre les premiers paragraphes (approximativement)
    paragraphs = content.split('\n\n')
    
    # Si trop de contenu, prendre seulement les 3 premiers paragraphes et ajouter "..."
    if len(paragraphs) > 3:
        text_frame = content_shape.text_frame
        text_frame.text = '\n'.join(paragraphs[:3])
        p = text_frame.add_paragraph()
        p.text = "..."
    else:
        content_shape.text = content


def integrer_avec_plateforme(
    contenu: str,
    plateforme: PlateformeIntegration,
    format_sortie: FormatOutput = FormatOutput.MARKDOWN,
    credentials: Optional[Dict] = None,
    metadata: Optional[Dict] = None,
) -> Dict:
    """
    Intègre un rapport avec une plateforme externe.
    
    Args:
        contenu: Le contenu du rapport
        plateforme: La plateforme d'intégration
        format_sortie: Le format de sortie à utiliser
        credentials: Informations d'authentification pour la plateforme
        metadata: Métadonnées pour le document
        
    Returns:
        Un dictionnaire avec les informations sur l'intégration (URL, ID, etc.)
    """
    # Note: Ceci est une implémentation fictive
    # Dans une implémentation réelle, vous devriez utiliser les API des plateformes
    
    result = {
        "success": True,
        "platform": plateforme.value,
        "message": f"Intégration fictive avec {plateforme.value} réussie"
    }
    
    if plateforme == PlateformeIntegration.NOTION:
        result["url"] = "https://notion.so/page-id"
        
    elif plateforme == PlateformeIntegration.GOOGLE_DOCS:
        result["url"] = "https://docs.google.com/document/d/doc-id"
        
    elif plateforme == PlateformeIntegration.SLACK:
        result["channel"] = metadata.get("channel", "general") if metadata else "general"
        
    elif plateforme == PlateformeIntegration.GITHUB:
        result["repo"] = metadata.get("repo", "user/repo") if metadata else "user/repo"
        result["path"] = metadata.get("path", "docs/report.md") if metadata else "docs/report.md"
        
    return result


if __name__ == "__main__":
    # Exemple d'utilisation
    test_markdown = """# Titre du rapport
    
## Section 1
    
Voici un paragraphe de test.
    
## Section 2
    
- Élément 1
- Élément 2
    
### Sous-section
    
Un autre paragraphe.
"""
    
    # Tester la conversion en HTML
    try:
        html_path = convertir_markdown_html(test_markdown, "test_rapport.html")
        print(f"HTML généré: {html_path}")
    except ImportError as e:
        print(f"Conversion HTML non disponible: {e}")
    
    # Tester la conversion en PDF
    try:
        pdf_path = convertir_markdown_pdf(test_markdown, "test_rapport.pdf")
        print(f"PDF généré: {pdf_path}")
    except ImportError as e:
        print(f"Conversion PDF non disponible: {e}")
